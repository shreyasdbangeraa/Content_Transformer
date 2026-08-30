import os
import re
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import Dict, Any, List
from app.config import settings

class PresentationGenerator:
    """Generates professional Microsoft PowerPoint (.pptx) presentations from slide JSON."""

    @staticmethod
    def render_pptx(deck_data: Dict[str, Any], output_filename: str) -> str:
        prs = Presentation()
        prs.slide_width = Inches(13.333) # 16:9 widescreen
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        # Theme Colors
        COLOR_BG_DARK = RGBColor(15, 23, 42)      # Slate 900
        COLOR_BG_LIGHT = RGBColor(248, 250, 252)  # Slate 50
        COLOR_ACCENT = RGBColor(14, 165, 233)     # Sky 500 / Cyan
        COLOR_TEXT_PRIMARY = RGBColor(15, 23, 42) # Slate 900
        COLOR_TEXT_MUTED = RGBColor(100, 116, 139)# Slate 500
        COLOR_CARD_BG = RGBColor(255, 255, 255)   # White

        slides: List[Dict[str, Any]] = deck_data.get("slides", [])
        deck_title = deck_data.get("deck_title", "Executive Presentation")

        # Slide 1: Title Slide
        slide_1 = prs.slides.add_slide(blank_layout)
        background = slide_1.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG_DARK

        # Title Box
        title_box = slide_1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11), Inches(3.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = deck_title.upper()
        p0.font.size = Pt(36)
        p0.font.bold = True
        p0.font.color.rgb = RGBColor(255, 255, 255)
        p0.font.name = 'Arial'

        p1 = tf.add_paragraph()
        p1.text = "conteX AI — Source-Grounded Executive Briefing"
        p1.font.size = Pt(18)
        p1.font.color.rgb = COLOR_ACCENT
        p1.font.name = 'Arial'
        p1.space_before = Pt(14)

        # Content Slides
        for s in slides:
            slide = prs.slides.add_slide(blank_layout)
            
            # Slide Header Banner
            header_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.33), Inches(1.4))
            htf = header_box.text_frame
            htf.word_wrap = True

            hp = htf.paragraphs[0]
            hp.text = s.get("title", "Key Briefing Point")
            hp.font.size = Pt(26)
            hp.font.bold = True
            hp.font.color.rgb = COLOR_TEXT_PRIMARY
            hp.font.name = 'Arial'

            sub = s.get("subtitle", "")
            if sub:
                hsub = htf.add_paragraph()
                hsub.text = sub
                hsub.font.size = Pt(14)
                hsub.font.color.rgb = COLOR_TEXT_MUTED
                hsub.font.name = 'Arial'
                hsub.space_before = Pt(4)

            # Bullet Content Box
            bullets = s.get("bullets", [])
            content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.5))
            ctf = content_box.text_frame
            ctf.word_wrap = True

            def _add_bullet_runs(paragraph, bullet_text: str):
                # Add bullet symbol
                bullet_symbol_run = paragraph.add_run()
                bullet_symbol_run.text = "•   "
                bullet_symbol_run.font.size = Pt(18)
                bullet_symbol_run.font.color.rgb = COLOR_ACCENT
                bullet_symbol_run.font.bold = True
                bullet_symbol_run.font.name = 'Arial'

                parts = re.split(r'(\*\*.*?\*\*|__.*?__|\*.*?\*|`.*?`)', bullet_text)
                for part in parts:
                    if not part:
                        continue
                    if (part.startswith('**') and part.endswith('**') and len(part) >= 4) or \
                       (part.startswith('__') and part.endswith('__') and len(part) >= 4):
                        run = paragraph.add_run()
                        run.text = part[2:-2]
                        run.font.bold = True
                        run.font.size = Pt(18)
                        run.font.color.rgb = COLOR_TEXT_PRIMARY
                        run.font.name = 'Arial'
                    elif part.startswith('*') and part.endswith('*') and len(part) >= 2:
                        run = paragraph.add_run()
                        run.text = part[1:-1]
                        run.font.italic = True
                        run.font.size = Pt(18)
                        run.font.color.rgb = COLOR_TEXT_PRIMARY
                        run.font.name = 'Arial'
                    else:
                        run = paragraph.add_run()
                        run.text = part
                        run.font.size = Pt(18)
                        run.font.color.rgb = COLOR_TEXT_PRIMARY
                        run.font.name = 'Arial'

            for idx, bullet in enumerate(bullets):
                bp = ctf.paragraphs[0] if idx == 0 else ctf.add_paragraph()
                _add_bullet_runs(bp, bullet)
                bp.space_before = Pt(12)

            # Speaker Notes
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = s.get("speaker_notes", "").replace("**", "")

        out_path = os.path.join(settings.EXPORT_DIR, output_filename)
        prs.save(out_path)
        return out_path
